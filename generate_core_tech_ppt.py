from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme
COLOR_PRIMARY = RGBColor(0x2B, 0x5F, 0xA6)  # #2B5FA6
COLOR_SECONDARY = RGBColor(0x4E, 0xCD, 0xC4)  # #4ECDC4
COLOR_ALERT = RGBColor(0xFF, 0x6B, 0x6B)  # #FF6B6B
COLOR_BG = RGBColor(0xF8, 0xFA, 0xFF)  # #F8FAFF
COLOR_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
COLOR_MUTED = RGBColor(0x78, 0x86, 0x99)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_CARD = RGBColor(0xEE, 0xF3, 0xFB)

FONT_NAME = "Microsoft YaHei"


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=16,
    bold=False,
    color=COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    font_name=FONT_NAME,
    margin=0.06,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    title,
    bullets,
    title_size=18,
    body_size=14,
    title_color=COLOR_PRIMARY,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = FONT_NAME
    r0.font.size = Pt(title_size)
    r0.font.bold = True
    r0.font.color.rgb = title_color

    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT_NAME
        p.font.size = Pt(body_size)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(2)

    return box


def add_rounded_card(slide, x, y, w, h, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.1):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    return shp


def add_title_block(slide, title, subtitle):
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.85)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_PRIMARY
    banner.line.fill.background()

    add_textbox(
        slide,
        0.45,
        0.1,
        7.8,
        0.38,
        title,
        size=24,
        bold=True,
        color=COLOR_WHITE,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        0.45,
        0.46,
        10.8,
        0.28,
        subtitle,
        size=12,
        bold=False,
        color=RGBColor(0xD9, 0xE8, 0xFF),
        align=PP_ALIGN.LEFT,
    )


def add_footer(slide, page_no, ref_text):
    add_textbox(
        slide,
        0.35,
        7.14,
        0.8,
        0.2,
        page_no,
        size=10,
        color=COLOR_MUTED,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        2.2,
        7.08,
        10.8,
        0.28,
        ref_text,
        size=9,
        color=COLOR_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_arrow(slide, x, y, w, h, color=COLOR_PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def add_vertical_arrow(slide, x, y, w, h, color=COLOR_PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def add_hline(slide, x1, y1, x2, y2, color=COLOR_PRIMARY, width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_circle(slide, x, y, w, h, fill=COLOR_SECONDARY, line=COLOR_PRIMARY):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = line
    c.line.width = Pt(1.0)
    return c


def make_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(
        slide,
        "核心技术（1/4）技术架构总览：微信小程序 × 事件驱动 × AI中台",
        "目标：一页讲清“采集-建档-评估-干预-随访”闭环，以及AI在医疗场景中的可控深度融合",
    )

    # Left 65% architecture area
    x0, y0, w0, h0 = 0.4, 1.0, 8.45, 5.9
    add_rounded_card(slide, x0, y0, w0, h0, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.4)
    add_textbox(slide, x0 + 0.18, y0 + 0.08, 3.4, 0.25, "专业技术架构图（示意）", size=12, bold=True, color=COLOR_PRIMARY)

    # Layer bands
    layer_h = 1.25
    layer_gap = 0.12
    layers = [
        ("前端层：微信小程序", COLOR_LIGHT_CARD),
        ("服务层：API网关 + 业务服务 + 事件总线", RGBColor(0xE8, 0xF8, 0xF6)),
        ("AI层：规则引擎 + RAG + 安全护栏", RGBColor(0xFF, 0xF4, 0xF4)),
        ("数据层：MySQL / Redis / 对象存储 / 审计日志", RGBColor(0xF3, 0xF5, 0xFA)),
    ]

    y = y0 + 0.4
    for name, c in layers:
        add_rounded_card(slide, x0 + 0.22, y, w0 - 0.44, layer_h, fill_color=c, line_color=RGBColor(0xC2, 0xCF, 0xE3), line_width=0.8)
        add_textbox(slide, x0 + 0.35, y + 0.04, 4.8, 0.24, name, size=11, bold=True, color=COLOR_PRIMARY)
        y += layer_h + layer_gap

    # Nodes in each layer
    add_rounded_card(slide, 0.95, 1.72, 1.15, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 1.0, 1.79, 1.05, 0.30, "疼痛标注", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 2.25, 1.72, 1.15, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 2.3, 1.79, 1.05, 0.30, "语音录入", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 3.55, 1.72, 1.15, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 3.6, 1.79, 1.05, 0.30, "AI问答", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 4.85, 1.72, 1.15, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 4.9, 1.79, 1.05, 0.30, "亲情号", size=10, align=PP_ALIGN.CENTER)

    add_rounded_card(slide, 1.05, 3.08, 1.5, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 1.1, 3.14, 1.4, 0.30, "API Gateway", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 2.75, 3.08, 1.55, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 2.8, 3.14, 1.45, 0.30, "记录/评估服务", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 4.5, 3.08, 1.4, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY)
    add_textbox(slide, 4.55, 3.14, 1.3, 0.30, "事件总线", size=10, align=PP_ALIGN.CENTER)

    add_rounded_card(slide, 1.0, 4.45, 1.7, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_ALERT)
    add_textbox(slide, 1.05, 4.51, 1.6, 0.30, "规则引擎", size=10, align=PP_ALIGN.CENTER, color=COLOR_ALERT, bold=True)
    add_rounded_card(slide, 2.95, 4.45, 1.55, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_ALERT)
    add_textbox(slide, 3.0, 4.51, 1.45, 0.30, "RAG检索", size=10, align=PP_ALIGN.CENTER, color=COLOR_ALERT, bold=True)
    add_rounded_card(slide, 4.75, 4.45, 1.45, 0.45, fill_color=COLOR_WHITE, line_color=COLOR_ALERT)
    add_textbox(slide, 4.8, 4.51, 1.35, 0.30, "安全护栏", size=10, align=PP_ALIGN.CENTER, color=COLOR_ALERT, bold=True)

    add_rounded_card(slide, 1.0, 5.82, 1.55, 0.45, fill_color=COLOR_WHITE, line_color=RGBColor(0x8A, 0x9C, 0xB8))
    add_textbox(slide, 1.05, 5.88, 1.45, 0.30, "MySQL", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 2.75, 5.82, 1.55, 0.45, fill_color=COLOR_WHITE, line_color=RGBColor(0x8A, 0x9C, 0xB8))
    add_textbox(slide, 2.8, 5.88, 1.45, 0.30, "Redis", size=10, align=PP_ALIGN.CENTER)
    add_rounded_card(slide, 4.5, 5.82, 1.9, 0.45, fill_color=COLOR_WHITE, line_color=RGBColor(0x8A, 0x9C, 0xB8))
    add_textbox(slide, 4.55, 5.88, 1.8, 0.30, "对象存储/审计", size=10, align=PP_ALIGN.CENTER)

    # Flow arrows
    add_vertical_arrow(slide, 6.8, 1.82, 0.35, 0.27, color=COLOR_PRIMARY)
    add_vertical_arrow(slide, 6.8, 3.18, 0.35, 0.27, color=COLOR_PRIMARY)
    add_vertical_arrow(slide, 6.8, 4.56, 0.35, 0.27, color=COLOR_PRIMARY)

    add_textbox(slide, 6.2, 2.22, 1.45, 0.22, "事件驱动", size=9, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.2, 3.58, 1.45, 0.22, "规则+AI协同", size=9, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.2, 4.95, 1.45, 0.22, "可追溯落库", size=9, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # Right side summary
    add_rounded_card(slide, 8.95, 1.0, 4.0, 5.9, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)

    add_bullets(
        slide,
        9.1,
        1.15,
        3.7,
        2.75,
        "核心技术亮点",
        [
            "事件驱动闭环：记录即触发评估、预警、推送与周报任务。",
            "医疗安全优先：规则/量表/状态机先决策，AI负责通俗化解释。",
            "AI深度融合：语音结构化、RAG证据检索、报告自动生成。",
        ],
        title_size=17,
        body_size=13,
    )

    add_bullets(
        slide,
        9.1,
        4.05,
        3.7,
        1.55,
        "落地形态（MVP）",
        [
            "微信小程序 + 云托管/Serverless：低运维、快迭代。",
            "Demo阶段优先开放：微信运动 + 订阅消息 + ASR。",
        ],
        title_size=15,
        body_size=12,
        title_color=COLOR_SECONDARY,
    )

    add_rounded_card(slide, 9.15, 5.78, 3.55, 0.9, fill_color=RGBColor(0xEF, 0xF8, 0xF7), line_color=COLOR_SECONDARY)
    add_textbox(
        slide,
        9.3,
        5.88,
        3.2,
        0.7,
        "工程口径更新：IMA 可作为扩展接口位；\n生产方案建议“本地可控RAG优先，外部检索补充”。",
        size=11,
        color=COLOR_TEXT,
    )

    add_footer(
        slide,
        "1/4",
        "参考：developers.weixin.qq.com；ima.qq.com；Tencent/WeKnora(README_CN)；J Med Internet Res 2020;22(9):e20283",
    )


def make_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(
        slide,
        "核心技术（2/4）创新一：零负担数据采集（破解高失访）",
        "目标：把“随访填表”改造成“10秒打卡 + AI自动结构化 + 一键确认”",
    )

    # Left panel
    add_rounded_card(slide, 0.4, 1.0, 6.35, 4.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)
    add_textbox(slide, 0.62, 1.12, 2.8, 0.26, "2D人体标注交互示意", size=13, bold=True, color=COLOR_PRIMARY)

    # Front and back mock body cards
    add_rounded_card(slide, 0.75, 1.58, 2.4, 3.3, fill_color=RGBColor(0xFB, 0xFD, 0xFF), line_color=RGBColor(0xC9, 0xD8, 0xEE), line_width=0.9)
    add_rounded_card(slide, 3.28, 1.58, 2.4, 3.3, fill_color=RGBColor(0xFB, 0xFD, 0xFF), line_color=RGBColor(0xC9, 0xD8, 0xEE), line_width=0.9)
    add_textbox(slide, 1.45, 1.7, 1.0, 0.2, "前视图", size=10, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.98, 1.7, 1.0, 0.2, "后视图", size=10, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # Body mock circles
    for cx, cy in [(1.8, 2.2), (1.8, 2.75), (1.8, 3.3), (1.35, 3.9), (2.25, 3.9)]:
        add_circle(slide, cx, cy, 0.28, 0.28, fill=RGBColor(0xD6, 0xEE, 0xFA), line=RGBColor(0xB3, 0xCB, 0xE6))
    for cx, cy in [(4.33, 2.2), (4.33, 2.75), (4.33, 3.3), (3.88, 3.9), (4.78, 3.9)]:
        add_circle(slide, cx, cy, 0.28, 0.28, fill=RGBColor(0xD6, 0xEE, 0xFA), line=RGBColor(0xB3, 0xCB, 0xE6))

    # Heat zones
    add_circle(slide, 1.7, 3.15, 0.5, 0.5, fill=RGBColor(0xFF, 0xA8, 0x9A), line=COLOR_ALERT)
    add_circle(slide, 4.2, 3.08, 0.65, 0.65, fill=RGBColor(0xFF, 0x7F, 0x7F), line=COLOR_ALERT)
    add_textbox(slide, 0.85, 4.98, 5.7, 0.35, "说明：点选区域 + 涂抹热区（颜色深浅=强度），只存region_id与评分，不存人体图片。", size=10, color=COLOR_MUTED)

    # Right panel
    add_rounded_card(slide, 6.95, 1.0, 6.0, 4.45, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)
    add_textbox(slide, 7.15, 1.12, 3.8, 0.26, "AI快速采集流程（语音→结构化→确认）", size=13, bold=True, color=COLOR_PRIMARY)

    steps = [
        "1) 用户语音：今天腰痛7分，走路加重，吃药后稍缓解",
        "2) ASR转写 + 字段抽取：部位/强度/诱因/用药/疗效",
        "3) 弹出确认卡片：确认 / 修改（大按钮，低阅读负担）",
        "4) 写入时间线并触发：风险评估 + 个体化宣教推送",
    ]
    sy = 1.56
    for i, s in enumerate(steps):
        card_color = RGBColor(0xF2, 0xF8, 0xFF) if i % 2 == 0 else RGBColor(0xF1, 0xFD, 0xFB)
        add_rounded_card(slide, 7.2, sy, 5.5, 0.72, fill_color=card_color, line_color=RGBColor(0xC3, 0xD4, 0xEA), line_width=0.8)
        add_textbox(slide, 7.35, sy + 0.15, 5.2, 0.5, s, size=11)
        if i < len(steps) - 1:
            add_vertical_arrow(slide, 9.7, sy + 0.72, 0.32, 0.21, color=COLOR_SECONDARY)
        sy += 0.96

    add_rounded_card(slide, 7.2, 4.62, 5.5, 0.72, fill_color=RGBColor(0xFF, 0xF4, 0xF4), line_color=COLOR_ALERT)
    add_textbox(
        slide,
        7.35,
        4.72,
        5.2,
        0.52,
        "触发补采集示例：步数骤降 / 服药后X小时 / 连续未记录 → 主动追问1句",
        size=11,
    )

    # Bottom comparison table
    add_rounded_card(slide, 0.4, 5.6, 12.55, 1.32, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.1)
    add_textbox(slide, 0.62, 5.68, 3.2, 0.24, "传统随访 vs 舒伴随访（关键差异）", size=12, bold=True, color=COLOR_PRIMARY)

    tbl = slide.shapes.add_table(5, 3, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.88)).table
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(4.8)
    tbl.columns[2].width = Inches(5.2)

    headers = ["关键数据", "传统问卷/人工随访", "舒伴 SoothPal（零负担采集）"]
    for c, t in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = t
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE9, 0xF0, 0xFB)

    rows = [
        ["疼痛部位/强度", "每日问卷 + 文字输入", "2D点选/涂抹 + 10秒打卡"],
        ["功能影响", "多题量表，易疲劳", "步数触发 + 1次点击"],
        ["用药疗效", "复诊回忆偏差", "服药后自动追踪 + AI确认卡片"],
        ["失访管理", "人工催访成本高", "状态机召回 + 亲情号联动"],
    ]

    for r in range(1, 5):
        for c in range(3):
            cell = tbl.cell(r, c)
            cell.text = rows[r - 1][c]
            if c == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xFC, 0xFA)

    for r in range(5):
        for c in range(3):
            tx = tbl.cell(r, c).text_frame
            for p in tx.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(10 if r > 0 else 10.5)
                    run.font.bold = r == 0
                    run.font.color.rgb = COLOR_TEXT
            tx.word_wrap = True

    add_footer(
        slide,
        "2/4",
        "参考：J Med Internet Res 2020;22(9):e20283（高失访问题）；团队调研与既往研究共识（30-60%区间用于问题界定）",
    )


def make_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(
        slide,
        "核心技术（3/4）创新二：危机预警与个性化推送（可解释、可配置、可审计）",
        "目标：突出医疗可靠性——规则底线 + 状态机个体化 + 分级响应 + 限频推送",
    )

    # Left logic panel
    add_rounded_card(slide, 0.4, 1.0, 7.1, 5.9, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)
    add_textbox(slide, 0.62, 1.12, 4.6, 0.25, "双路径预警逻辑图", size=13, bold=True, color=COLOR_PRIMARY)

    # Path A
    add_rounded_card(slide, 0.72, 1.55, 3.0, 2.12, fill_color=RGBColor(0xFF, 0xF2, 0xF2), line_color=COLOR_ALERT, line_width=1.0)
    add_textbox(slide, 0.9, 1.68, 2.6, 0.28, "Path A：静态阈值/突变（急性）", size=11, bold=True, color=COLOR_ALERT)
    add_textbox(slide, 0.95, 2.02, 2.45, 1.45, "• NRS ≥ 8\n• 当日较7日均值突增 ≥ 3 且当日 ≥ 6\n• 优先级最高，立即触发", size=10.5)

    # Path B
    add_rounded_card(slide, 3.98, 1.55, 3.2, 2.12, fill_color=RGBColor(0xF0, 0xFA, 0xF9), line_color=COLOR_SECONDARY, line_width=1.0)
    add_textbox(slide, 4.16, 1.68, 2.8, 0.28, "Path B：治疗周期状态机（慢性）", size=11, bold=True, color=RGBColor(0x17, 0x9D, 0x95))
    add_textbox(slide, 4.2, 2.02, 2.7, 1.45, "• 术后第1周 / 康复期 / 药物敏感期\n• 阈值动态调整 + 连续性判定\n• 识别持续高痛与依从性下降", size=10.5)

    add_vertical_arrow(slide, 2.95, 3.78, 0.32, 0.27, color=COLOR_PRIMARY)
    add_vertical_arrow(slide, 5.45, 3.78, 0.32, 0.27, color=COLOR_PRIMARY)

    add_rounded_card(slide, 2.2, 4.15, 3.5, 1.05, fill_color=RGBColor(0xEC, 0xF2, 0xFF), line_color=COLOR_PRIMARY)
    add_textbox(slide, 2.35, 4.31, 3.2, 0.28, "Risk Level 1/2/3", size=13, bold=True, align=PP_ALIGN.CENTER, color=COLOR_PRIMARY)
    add_textbox(slide, 2.35, 4.62, 3.2, 0.30, "Action Dispatcher：通知 / 报告 / 宣教", size=10.5, align=PP_ALIGN.CENTER)

    add_vertical_arrow(slide, 3.78, 5.24, 0.32, 0.25, color=COLOR_PRIMARY)
    add_rounded_card(slide, 1.85, 5.52, 4.2, 1.08, fill_color=RGBColor(0xF8, 0xFC, 0xFF), line_color=RGBColor(0xB7, 0xC7, 0xDE))
    add_textbox(
        slide,
        2.02,
        5.68,
        3.9,
        0.78,
        "执行方式：触发器/云函数实时计算，\n输出携带 rule_id + reason，支持审计复盘。",
        size=11,
        align=PP_ALIGN.CENTER,
    )

    # Right panel
    add_rounded_card(slide, 7.7, 1.0, 5.25, 5.9, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)
    add_textbox(slide, 7.92, 1.12, 2.9, 0.25, "分级响应（示例）", size=13, bold=True, color=COLOR_PRIMARY)

    rows = [
        ("Level 1 高危", "NRS≥8 或突增>3", "强提醒 + 就医模板 + 紧急摘要 + (授权)亲情号通知", COLOR_ALERT),
        ("Level 2 中危", "5-7分持续≥3天/阶段异常", "建议联系医生 + 并发症/康复注意事项推送", RGBColor(0xE6, 0x8A, 0x00)),
        ("Level 3 依从性", "超过X天未记录且未暂停", "温和召回 + 24h延迟队列 + 亲情号兜底", RGBColor(0x4C, 0x8D, 0xD6)),
    ]

    y = 1.5
    for title, cond, act, c in rows:
        add_rounded_card(slide, 7.95, y, 4.75, 1.22, fill_color=RGBColor(0xFB, 0xFD, 0xFF), line_color=c, line_width=1.0)
        add_textbox(slide, 8.1, y + 0.08, 1.75, 0.2, title, size=11, bold=True, color=c)
        add_textbox(slide, 8.1, y + 0.33, 4.45, 0.24, f"触发：{cond}", size=10)
        add_textbox(slide, 8.1, y + 0.58, 4.45, 0.5, f"动作：{act}", size=10)
        y += 1.35

    add_rounded_card(slide, 7.95, 5.6, 4.75, 1.0, fill_color=RGBColor(0xF1, 0xFD, 0xFB), line_color=COLOR_SECONDARY)
    add_textbox(
        slide,
        8.1,
        5.72,
        4.45,
        0.74,
        "推送混合策略：\n日常按画像匹配，危机消息置顶覆盖常规科普；\n全程去重 + 限频，减少打扰并提升留存。",
        size=10.5,
    )

    add_footer(
        slide,
        "3/4",
        "参考：支瑞聪等,2020；朱南希等,2022；团队规则引擎设计稿（可配置阈值+状态机）",
    )


def make_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_block(
        slide,
        "核心技术（4/4）创新三：AI深度融合（RAG+脱敏检索+安全护栏）与可行性",
        "目标：给出可工程落地且可答辩的“隐私保护+可引用输出”方案（Demo→试点）",
    )

    # Top sequence area
    add_rounded_card(slide, 0.4, 1.0, 12.55, 3.5, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)
    add_textbox(slide, 0.62, 1.12, 4.5, 0.25, "RAG/Agent 时序流程（上：泳道图）", size=13, bold=True, color=COLOR_PRIMARY)

    # Swimlanes
    lane_y = [1.5, 2.25, 3.0]
    lane_titles = ["小程序Agent", "AI网关（后端）", "知识检索层（本地RAG优先）"]
    for i, y in enumerate(lane_y):
        add_rounded_card(slide, 0.7, y, 12.0, 0.58, fill_color=RGBColor(0xFA, 0xFC, 0xFF), line_color=RGBColor(0xD3, 0xDF, 0xEF), line_width=0.8)
        add_textbox(slide, 0.85, y + 0.16, 2.2, 0.24, lane_titles[i], size=10.5, bold=True, color=COLOR_PRIMARY)

    # Sequence boxes
    add_rounded_card(slide, 2.9, 1.62, 1.7, 0.3, fill_color=RGBColor(0xE9, 0xF0, 0xFB), line_color=COLOR_PRIMARY, line_width=0.8)
    add_textbox(slide, 2.98, 1.68, 1.55, 0.18, "接收用户提问", size=9.5, align=PP_ALIGN.CENTER)

    add_rounded_card(slide, 4.95, 2.37, 1.9, 0.3, fill_color=RGBColor(0xFF, 0xF3, 0xF3), line_color=COLOR_ALERT, line_width=0.8)
    add_textbox(slide, 5.03, 2.43, 1.75, 0.18, "脱敏/最小必要Query", size=9.5, align=PP_ALIGN.CENTER, color=COLOR_ALERT, bold=True)

    add_rounded_card(slide, 7.25, 3.12, 2.05, 0.3, fill_color=RGBColor(0xF0, 0xFB, 0xFA), line_color=COLOR_SECONDARY, line_width=0.8)
    add_textbox(slide, 7.33, 3.18, 1.9, 0.18, "检索证据片段（带来源）", size=9.5, align=PP_ALIGN.CENTER)

    add_rounded_card(slide, 9.65, 2.37, 2.15, 0.3, fill_color=RGBColor(0xFF, 0xF3, 0xF3), line_color=COLOR_ALERT, line_width=0.8)
    add_textbox(slide, 9.73, 2.43, 2.0, 0.18, "安全护栏+通俗化生成", size=9.5, align=PP_ALIGN.CENTER, color=COLOR_ALERT, bold=True)

    add_rounded_card(slide, 8.25, 1.62, 2.05, 0.3, fill_color=RGBColor(0xE9, 0xF0, 0xFB), line_color=COLOR_PRIMARY, line_width=0.8)
    add_textbox(slide, 8.33, 1.68, 1.9, 0.18, "返回短句答案+引用", size=9.5, align=PP_ALIGN.CENTER)

    add_arrow(slide, 4.62, 1.67, 0.28, 0.18, color=COLOR_PRIMARY)
    add_arrow(slide, 6.88, 2.42, 0.3, 0.18, color=COLOR_PRIMARY)
    add_arrow(slide, 9.35, 3.17, 0.24, 0.18, color=COLOR_PRIMARY)
    add_arrow(slide, 10.7, 1.67, 0.38, 0.18, color=COLOR_PRIMARY)

    # Privacy shields
    shield1 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(3.05), Inches(2.9), Inches(0.52), Inches(0.35))
    shield1.fill.solid()
    shield1.fill.fore_color.rgb = COLOR_ALERT
    shield1.line.fill.background()
    add_textbox(slide, 3.1, 2.98, 0.42, 0.16, "隐私", size=8.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    shield2 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(11.0), Inches(2.9), Inches(0.52), Inches(0.35))
    shield2.fill.solid()
    shield2.fill.fore_color.rgb = COLOR_ALERT
    shield2.line.fill.background()
    add_textbox(slide, 11.05, 2.98, 0.42, 0.16, "护栏", size=8.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        0.8,
        3.65,
        11.9,
        0.58,
        "隐私不出域：用户画像、原始病历附件不外发；外部检索仅接收脱敏后的最小必要Query。",
        size=10.5,
        color=COLOR_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Bottom cards
    add_rounded_card(slide, 0.4, 4.65, 12.55, 2.23, fill_color=COLOR_WHITE, line_color=COLOR_PRIMARY, line_width=1.2)

    card_w = 3.95
    gaps = 0.28
    x_start = 0.7

    # Card 1
    add_rounded_card(slide, x_start, 4.92, card_w, 1.75, fill_color=RGBColor(0xF1, 0xF7, 0xFF), line_color=COLOR_PRIMARY, line_width=0.9)
    add_textbox(slide, x_start + 0.14, 5.04, 3.65, 0.25, "知识库资产化（专业壁垒）", size=12, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, x_start + 0.16, 5.33, 3.55, 1.2, "• 指南/共识/循证/医生科普结构化切片\n• 向量检索 + 版本管理 + 专家审核\n• 输出带引用，可追溯可质控", size=10.5)

    # Card 2
    x2 = x_start + card_w + gaps
    add_rounded_card(slide, x2, 4.92, card_w, 1.75, fill_color=RGBColor(0xF2, 0xFC, 0xFA), line_color=COLOR_SECONDARY, line_width=0.9)
    add_textbox(slide, x2 + 0.14, 5.04, 3.65, 0.25, "安全合规护栏（医疗可控）", size=12, bold=True, color=RGBColor(0x1E, 0x98, 0x90))
    add_textbox(slide, x2 + 0.16, 5.33, 3.55, 1.2, "• 先规则后生成：风险引擎判定“该不该做”\n• AI负责“怎么说/怎么教”，禁止处方剂量建议\n• 同意台账 + 权限控制 + 审计留痕", size=10.5)

    # Card 3
    x3 = x2 + card_w + gaps
    add_rounded_card(slide, x3, 4.92, card_w, 1.75, fill_color=RGBColor(0xFF, 0xF5, 0xF5), line_color=COLOR_ALERT, line_width=0.9)
    add_textbox(slide, x3 + 0.14, 5.04, 3.65, 0.25, "成熟度与扩展（Demo→试点）", size=12, bold=True, color=COLOR_ALERT)
    add_textbox(slide, x3 + 0.16, 5.33, 3.55, 1.2, "• Demo：语音结构化、RAG引用回答、分级预警、亲情号通知\n• 扩展：更多可穿戴与健康平台数据接入\n• 可选：IMA Skill接口位（以官方能力开放节奏为准）", size=10.5)

    add_footer(
        slide,
        "4/4",
        "参考：ima.qq.com/agent-interface；Tencent/WeKnora(README_CN, API文档)；彭进亚等2016；朱南希等2022；李冬2019",
    )


def build_ppt(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_slide_1(prs)
    make_slide_2(prs)
    make_slide_3(prs)
    make_slide_4(prs)

    prs.save(output_path)


if __name__ == "__main__":
    out = "d:/project/pain_alleviation/SoothPal_核心技术_16x9_最终版.pptx"
    build_ppt(out)
    print(f"Generated: {out}")
