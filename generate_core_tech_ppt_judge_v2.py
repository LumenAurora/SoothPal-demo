from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color system
PRIMARY = RGBColor(0x2B, 0x5F, 0xA6)
SECONDARY = RGBColor(0x25, 0xB5, 0xA9)
ALERT = RGBColor(0xE8, 0x56, 0x56)
BG = RGBColor(0xF8, 0xFA, 0xFF)
TEXT = RGBColor(0x1F, 0x2D, 0x3D)
MUTED = RGBColor(0x6E, 0x7F, 0x97)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_A = RGBColor(0xEE, 0xF4, 0xFF)
CARD_B = RGBColor(0xEE, 0xFB, 0xF9)
CARD_C = RGBColor(0xFF, 0xF3, 0xF3)

FONT = "Microsoft YaHei"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def txt(slide, x, y, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, title, items, title_color=PRIMARY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = FONT
    r0.font.bold = True
    r0.font.size = Pt(14)
    r0.font.color.rgb = title_color

    for it in items:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT
        p.space_before = Pt(2)

    return box


def card(slide, x, y, w, h, fill=WHITE, line=PRIMARY, lw=1.0, rounded=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    return shp


def arrow_right(slide, x, y, w=0.34, h=0.18, color=PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def arrow_down(slide, x, y, w=0.26, h=0.2, color=PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def title_bar(slide, title, sub):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.82))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()

    txt(slide, 0.45, 0.09, 9.5, 0.32, title, size=24, bold=True, color=WHITE)
    txt(slide, 0.45, 0.46, 12.2, 0.22, sub, size=11, color=RGBColor(0xD9, 0xE7, 0xFB))


def footer(slide, page, ref):
    txt(slide, 0.35, 7.15, 0.8, 0.2, page, size=10, color=MUTED)
    txt(slide, 2.0, 7.08, 11.0, 0.28, ref, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "核心技术（1/4）总体架构：低负担采集 × 可解释预警 × AI中台", "评委关注点：可行性、专业性、可复制性三者同时成立")

    # Left architecture area
    card(s, 0.4, 1.05, 8.55, 5.82, fill=WHITE, line=PRIMARY, lw=1.2)
    txt(s, 0.62, 1.14, 3.2, 0.24, "技术闭环（采集-评估-干预-随访）", size=13, bold=True, color=PRIMARY)

    labels = [
        ("零负担采集", "2D点选/语音打卡\n自动补采集", CARD_A),
        ("数据治理", "结构化档案\n时间线事件", CARD_A),
        ("风险评估", "阈值+突变+状态机\nrule_id可审计", CARD_C),
        ("AI引擎", "RAG证据检索\n通俗化生成", CARD_B),
        ("触达闭环", "分级提醒/周报\n亲情号联动", CARD_B),
    ]

    x0 = 0.72
    y0 = 2.05
    w = 1.45
    h = 1.35
    gap = 0.22
    for i, (t, d, c) in enumerate(labels):
        x = x0 + i * (w + gap)
        card(s, x, y0, w, h, fill=c, line=RGBColor(0xC3, 0xD2, 0xE6), lw=0.8)
        txt(s, x + 0.08, y0 + 0.12, w - 0.16, 0.3, t, size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        txt(s, x + 0.08, y0 + 0.48, w - 0.16, 0.65, d, size=10.2, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_right(s, x + w + 0.03, y0 + 0.56, 0.16, 0.2)

    # Tech stack strip
    card(s, 0.62, 4.08, 8.1, 2.55, fill=RGBColor(0xF7, 0xFA, 0xFF), line=RGBColor(0xC5, 0xD3, 0xE6), lw=0.8)
    txt(s, 0.82, 4.2, 2.5, 0.24, "推荐技术栈（MVP可落地）", size=12, bold=True, color=PRIMARY)

    chips = [
        "微信小程序", "云托管/Serverless", "MySQL", "Redis", "对象存储", "ASR", "RAG知识库", "订阅消息", "审计日志"
    ]
    cx, cy = 0.85, 4.55
    for i, c in enumerate(chips):
        cw = 1.2 if len(c) <= 6 else 1.65
        if cx + cw > 8.35:
            cx = 0.85
            cy += 0.52
        card(s, cx, cy, cw, 0.38, fill=WHITE, line=RGBColor(0xAF, 0xC1, 0xD9), lw=0.8)
        txt(s, cx + 0.05, cy + 0.09, cw - 0.1, 0.2, c, size=9.8, align=PP_ALIGN.CENTER)
        cx += cw + 0.14

    # Right value panel
    card(s, 9.15, 1.05, 3.78, 5.82, fill=WHITE, line=PRIMARY, lw=1.2)
    bullet_box(
        s, 9.32, 1.25, 3.42, 2.0, "评委视角价值", [
            "可行：团队现有能力可在赛期内完成Demo闭环",
            "可信：规则先决策，AI负责解释与教育",
            "可扩：B2B2C场景下可复制到更多慢病管理",
        ]
    )
    bullet_box(
        s, 9.32, 3.35, 3.42, 1.55, "差异化壁垒", [
            "低负担采集带来更高随访完成率",
            "可解释预警与审计链条适配医疗场景",
        ], title_color=SECONDARY
    )
    card(s, 9.32, 5.0, 3.42, 1.6, fill=RGBColor(0xF1, 0xF8, 0xFF), line=RGBColor(0xAC, 0xBF, 0xDD), lw=0.8)
    txt(s, 9.48, 5.17, 3.1, 1.2, "一句话总结：\n不是“聊天型健康App”，\n而是“可追溯、可干预、可规模化”的数字疼痛管理底座。", size=11)

    footer(s, "1/4", "参考：developers.weixin.qq.com；ima.qq.com；J Med Internet Res 2020;22(9):e20283")


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "核心技术（2/4）创新一：零负担数据采集（10秒完成关键记录）", "评委关注点：降低失访、降低人工随访成本、提高数据连续性")

    # Two balanced panels
    card(s, 0.4, 1.05, 6.25, 5.2, fill=WHITE, line=PRIMARY, lw=1.2)
    card(s, 6.85, 1.05, 6.1, 5.2, fill=WHITE, line=PRIMARY, lw=1.2)

    txt(s, 0.62, 1.15, 3.2, 0.25, "交互入口：2D人体点选/涂抹", size=13, bold=True, color=PRIMARY)
    # phone mock
    card(s, 1.0, 1.62, 2.18, 3.85, fill=RGBColor(0xFB, 0xFD, 0xFF), line=RGBColor(0xBB, 0xCB, 0xE1), lw=0.9)
    txt(s, 1.15, 1.75, 1.85, 0.2, "前视图", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    card(s, 3.45, 1.62, 2.18, 3.85, fill=RGBColor(0xFB, 0xFD, 0xFF), line=RGBColor(0xBB, 0xCB, 0xE1), lw=0.9)
    txt(s, 3.6, 1.75, 1.85, 0.2, "后视图", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # simple heat points
    for x, y in [(1.95, 2.35), (1.95, 2.9), (1.95, 3.45), (1.5, 4.2), (2.4, 4.2), (4.4, 2.35), (4.4, 2.9), (4.4, 3.45), (3.95, 4.2), (4.85, 4.2)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.28), Inches(0.28))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xC6, 0xE2, 0xF5)
        c.line.color.rgb = RGBColor(0x9B, 0xC2, 0xE0)
    hot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.72), Inches(3.2), Inches(0.55), Inches(0.55))
    hot.fill.solid(); hot.fill.fore_color.rgb = RGBColor(0xF6, 0x9D, 0x94); hot.line.color.rgb = ALERT
    hot2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.1), Inches(3.1), Inches(0.68), Inches(0.68))
    hot2.fill.solid(); hot2.fill.fore_color.rgb = RGBColor(0xF2, 0x7A, 0x7A); hot2.line.color.rgb = ALERT

    txt(s, 0.75, 5.58, 5.8, 0.5, "关键点：\n点选区域 + 强度滑条 + 一键保存，不要求长文本输入。", size=10.5, color=MUTED)

    txt(s, 7.07, 1.15, 3.8, 0.25, "AI结构化采集流程", size=13, bold=True, color=PRIMARY)
    steps = [
        "1) 用户一句话描述疼痛与用药",
        "2) ASR转写 + 字段抽取（部位/强度/诱因/疗效）",
        "3) 弹出确认卡片（确认/修改）",
        "4) 自动写入时间线并触发风险评估",
    ]
    yy = 1.65
    for i, st in enumerate(steps):
        fill = CARD_A if i % 2 == 0 else CARD_B
        card(s, 7.12, yy, 5.55, 0.85, fill=fill, line=RGBColor(0xBE, 0xD0, 0xE5), lw=0.8)
        txt(s, 7.28, yy + 0.24, 5.2, 0.42, st, size=11.3)
        if i < 3:
            arrow_down(s, 9.7, yy + 0.86, 0.26, 0.2, color=SECONDARY)
        yy += 1.05

    card(s, 7.12, 5.1, 5.55, 0.95, fill=RGBColor(0xFF, 0xF6, 0xF6), line=ALERT, lw=0.9)
    txt(s, 7.28, 5.34, 5.2, 0.52, "触发式补采集：步数骤降 / 服药后X小时 / 连续未记录 → 主动追问1句", size=11)

    # bottom KPI strip
    card(s, 0.4, 6.35, 12.55, 0.63, fill=WHITE, line=PRIMARY, lw=1.1)
    card(s, 0.62, 6.45, 3.95, 0.42, fill=CARD_A, line=RGBColor(0xB8, 0xCB, 0xE3), lw=0.8)
    card(s, 4.8, 6.45, 3.95, 0.42, fill=CARD_B, line=RGBColor(0xA8, 0xD9, 0xD3), lw=0.8)
    card(s, 8.98, 6.45, 3.75, 0.42, fill=CARD_C, line=RGBColor(0xE8, 0xB8, 0xB8), lw=0.8)
    txt(s, 0.8, 6.55, 3.6, 0.22, "采集耗时：目标≤10秒", size=10.8, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    txt(s, 4.98, 6.55, 3.6, 0.22, "字段结构化：自动完成", size=10.8, bold=True, color=RGBColor(0x1B, 0x9A, 0x92), align=PP_ALIGN.CENTER)
    txt(s, 9.15, 6.55, 3.45, 0.22, "失访召回：状态机触发", size=10.8, bold=True, color=ALERT, align=PP_ALIGN.CENTER)

    footer(s, "2/4", "参考：J Med Internet Res 2020;22(9):e20283；团队一线访谈与问卷调研")


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "核心技术（3/4）创新二：可解释预警与个性化推送", "评委关注点：医疗安全可靠、策略可配置、输出可审计")

    card(s, 0.4, 1.05, 7.2, 5.85, fill=WHITE, line=PRIMARY, lw=1.2)
    card(s, 7.82, 1.05, 5.13, 5.85, fill=WHITE, line=PRIMARY, lw=1.2)

    txt(s, 0.62, 1.15, 3.2, 0.25, "双路径风险识别", size=13, bold=True, color=PRIMARY)
    card(s, 0.8, 1.62, 3.15, 2.08, fill=RGBColor(0xFF, 0xF4, 0xF4), line=ALERT, lw=0.9)
    bullet_box(s, 0.92, 1.78, 2.9, 1.8, "Path A：静态阈值/突变", [
        "NRS ≥ 8 直接触发高危",
        "较7日均值突增 ≥ 3 且当日≥6",
        "用于急性风险快速兜底",
    ], title_color=ALERT)

    card(s, 4.1, 1.62, 3.15, 2.08, fill=RGBColor(0xF0, 0xFC, 0xFA), line=SECONDARY, lw=0.9)
    bullet_box(s, 4.22, 1.78, 2.9, 1.8, "Path B：治疗周期状态机", [
        "术后期/康复期/药物敏感期",
        "阈值动态调整 + 连续性判定",
        "用于慢性风险和依从性识别",
    ], title_color=RGBColor(0x1E, 0x98, 0x90))

    arrow_down(s, 2.2, 3.82, 0.3, 0.25)
    arrow_down(s, 5.52, 3.82, 0.3, 0.25)

    card(s, 2.1, 4.2, 3.8, 1.02, fill=RGBColor(0xEC, 0xF2, 0xFF), line=PRIMARY, lw=1.0)
    txt(s, 2.25, 4.43, 3.5, 0.24, "Risk Level 1 / 2 / 3", size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    txt(s, 2.25, 4.72, 3.5, 0.26, "Action Dispatcher：通知 / 报告 / 宣教", size=11.2, align=PP_ALIGN.CENTER)

    arrow_down(s, 3.86, 5.24, 0.28, 0.2)
    card(s, 1.4, 5.52, 5.18, 1.15, fill=RGBColor(0xF8, 0xFB, 0xFF), line=RGBColor(0xBC, 0xCB, 0xE1), lw=0.8)
    txt(s, 1.62, 5.8, 4.75, 0.62, "实时监听：触发器/云函数计算；\n审计输出：rule_id + reason + event_time，便于复盘。", size=11, align=PP_ALIGN.CENTER)

    txt(s, 8.03, 1.15, 2.8, 0.25, "分级响应与推送策略", size=13, bold=True, color=PRIMARY)

    blocks = [
        ("Level 1 高危", "NRS≥8或突增>3", "强提醒 + 就医建议模板 + 紧急摘要", ALERT),
        ("Level 2 中危", "5-7分持续≥3天", "建议联系医生 + 并发症注意事项推送", RGBColor(0xD9, 0x8B, 0x16)),
        ("Level 3 依从性", "超过X天未记录", "温和召回 + 24h延迟队列 + 亲情号提醒", RGBColor(0x4B, 0x89, 0xD6)),
    ]
    y = 1.6
    for t, c1, c2, lc in blocks:
        card(s, 8.06, y, 4.68, 1.35, fill=WHITE, line=lc, lw=1.0)
        txt(s, 8.22, y + 0.13, 1.6, 0.25, t, size=12, bold=True, color=lc)
        txt(s, 8.22, y + 0.44, 4.35, 0.24, f"触发：{c1}", size=10.8)
        txt(s, 8.22, y + 0.72, 4.35, 0.42, f"动作：{c2}", size=10.8)
        y += 1.55

    card(s, 8.06, 6.25, 4.68, 0.52, fill=CARD_B, line=SECONDARY, lw=0.8)
    txt(s, 8.22, 6.37, 4.35, 0.28, "推送规则：危机信息置顶覆盖常规科普；全程去重 + 限频。", size=10.8)

    footer(s, "3/4", "参考：支瑞聪等,2020；朱南希等,2022；团队预警规则设计")


def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_bar(s, "核心技术（4/4）创新三：AI深度融合与技术壁垒", "评委关注点：AI不是噱头，而是可控可追溯的生产力中枢")

    card(s, 0.4, 1.05, 12.55, 3.35, fill=WHITE, line=PRIMARY, lw=1.2)
    txt(s, 0.62, 1.15, 3.7, 0.25, "AI闭环流程（上：时序）", size=13, bold=True, color=PRIMARY)

    # horizontal flow
    steps = [
        ("用户提问", CARD_A, PRIMARY),
        ("AI网关脱敏", CARD_C, ALERT),
        ("双引擎检索\n(IMA知识库+平台知识库)", CARD_B, SECONDARY),
        ("安全护栏与通俗化", CARD_C, ALERT),
        ("返回短句答案+引用", CARD_A, PRIMARY),
    ]
    sx, sy = 0.86, 2.12
    sw, sh, sg = 2.2, 1.38, 0.2
    for i, (name, fc, lc) in enumerate(steps):
        card(s, sx + i * (sw + sg), sy, sw, sh, fill=fc, line=lc, lw=0.9)
        txt(s, sx + i * (sw + sg) + 0.1, sy + 0.42, sw - 0.2, 0.55, name, size=11.3, bold=True, color=lc, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_right(s, sx + sw + i * (sw + sg) + 0.03, sy + 0.56, 0.14, 0.22)

    card(s, 0.8, 3.65, 12.0, 0.55, fill=RGBColor(0xF7, 0xFB, 0xFF), line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.8)
    txt(s, 1.0, 3.82, 11.6, 0.25, "隐私原则：用户画像与原始病历不外发；外部能力仅接收脱敏后的最小必要Query。", size=10.8, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom 3 cards + summary
    card(s, 0.4, 4.55, 12.55, 2.35, fill=WHITE, line=PRIMARY, lw=1.2)

    cw, ch = 3.9, 1.65
    xs = [0.7, 4.72, 8.74]
    fills = [CARD_A, CARD_B, CARD_C]
    titles = [
        "壁垒1：知识资产化",
        "壁垒2：安全可控",
        "壁垒3：规模化能力",
    ]
    contents = [
        ["指南/共识/循证材料结构化沉淀", "输出带引用，可追溯可质控", "专家审核保障医学可靠性"],
        ["规则先决策，AI做解释与教育", "高危场景使用固定就医模板", "权限、同意、审计全链条留痕"],
        ["云函数调度支持高并发低运维", "可扩展接入更多可穿戴平台", "B2B2C路径支撑商业化推广"],
    ]
    title_colors = [PRIMARY, RGBColor(0x1E, 0x98, 0x90), ALERT]

    for i in range(3):
        card(s, xs[i], 4.82, cw, ch, fill=fills[i], line=title_colors[i], lw=0.9)
        bullet_box(s, xs[i] + 0.08, 4.95, cw - 0.16, ch - 0.14, titles[i], contents[i], title_color=title_colors[i])

    card(s, 0.7, 6.55, 12.1, 0.25, fill=RGBColor(0xEC, 0xF2, 0xFF), line=RGBColor(0xB3, 0xC6, 0xE0), lw=0.7)
    txt(s, 0.9, 6.57, 11.7, 0.18, "结论：本项目具备“技术可行 + 医疗可信 + 商业可扩”的三重成立条件。", size=10.6, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    footer(s, "4/4", "参考：ima.qq.com；developers.weixin.qq.com；彭进亚等,2016；朱南希等,2022；李冬,2019")


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)

    prs.save(path)


if __name__ == "__main__":
    output = "d:/project/pain_alleviation/SoothPal_核心技术_评委优化版.pptx"
    build(output)
    print(output)
