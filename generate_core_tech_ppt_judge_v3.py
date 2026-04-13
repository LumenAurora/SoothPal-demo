from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme
PRIMARY = RGBColor(0x2B, 0x5F, 0xA6)
ACCENT = RGBColor(0x22, 0xB8, 0xA9)
ALERT = RGBColor(0xE6, 0x5A, 0x5A)
BG = RGBColor(0xF8, 0xFA, 0xFF)
TEXT = RGBColor(0x1E, 0x2A, 0x3A)
MUTED = RGBColor(0x6B, 0x7D, 0x95)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BLUE = RGBColor(0xEE, 0xF4, 0xFF)
CARD_TEAL = RGBColor(0xEE, 0xFB, 0xF9)
CARD_RED = RGBColor(0xFF, 0xF4, 0xF4)
CARD_GRAY = RGBColor(0xF3, 0xF6, 0xFC)

FONT = "Microsoft YaHei"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(
    slide,
    x,
    y,
    w,
    h,
    content,
    size=14,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, title, lines, title_color=PRIMARY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.04)

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = FONT
    r0.font.size = Pt(13)
    r0.font.bold = True
    r0.font.color.rgb = title_color

    for line in lines:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.space_before = Pt(2)


def add_card(slide, x, y, w, h, fill=WHITE, line=PRIMARY, lw=0.9, rounded=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    return shp


def add_topbar(slide, title, subtitle):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.88))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    add_text(slide, 0.48, 0.09, 11.9, 0.34, title, size=22, bold=True, color=WHITE)
    add_text(slide, 0.48, 0.52, 12.1, 0.2, subtitle, size=10, color=RGBColor(0xD9, 0xE7, 0xFB))


def add_footer(slide, page, refs):
    add_text(slide, 0.36, 7.15, 0.8, 0.2, page, size=10, color=MUTED)
    add_text(slide, 2.1, 7.08, 10.9, 0.25, refs, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def arrow_right(slide, x, y, w=0.25, h=0.18, color=PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def arrow_down(slide, x, y, w=0.22, h=0.18, color=PRIMARY):
    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_topbar(
        s,
        "核心技术（1/4）技术闭环架构：从记录到干预",
        "舒伴 SoothPal 以“低负担采集 + 可解释预警 + AI赋能教育”构建慢性疼痛管理底座",
    )

    # Main areas
    add_card(s, 0.55, 1.1, 8.75, 3.45, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 0.55, 4.75, 8.75, 2.1, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 9.5, 1.1, 3.25, 5.75, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)

    add_text(s, 0.75, 1.22, 3.8, 0.24, "核心流程：采集 → 建档 → 评估 → 干预 → 随访", size=12.5, bold=True, color=PRIMARY)

    steps = [
        ("零负担采集", "2D点选/语音打卡", CARD_BLUE),
        ("数据建档", "时间线+用户画像", CARD_BLUE),
        ("风险评估", "阈值+突变+状态机", CARD_RED),
        ("AI生成", "检索增强+通俗表达", CARD_TEAL),
        ("闭环随访", "分级提醒+周报", CARD_TEAL),
    ]
    sx, sy = 0.8, 2.0
    sw, sh, sg = 1.55, 1.45, 0.18
    for i, (t, d, c) in enumerate(steps):
        x = sx + i * (sw + sg)
        add_card(s, x, sy, sw, sh, fill=c, line=RGBColor(0xB4, 0xC6, 0xDE), lw=0.8)
        add_text(s, x + 0.08, sy + 0.18, sw - 0.16, 0.28, t, size=11.8, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.08, sy + 0.58, sw - 0.16, 0.58, d, size=10.2, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow_right(s, x + sw + 0.03, sy + 0.62, 0.12, 0.18)

    add_text(s, 0.75, 4.87, 3.3, 0.24, "关键技术栈（MVP可直接落地）", size=12.5, bold=True, color=PRIMARY)

    groups = [
        ("前端与交互", ["微信小程序", "2D疼痛标注", "语音录入(ASR)"]),
        ("服务与数据", ["云托管/Serverless", "MySQL + Redis", "对象存储 + 审计日志"]),
        ("智能与触达", ["RAG知识引擎", "规则预警引擎", "订阅消息/亲情号通知"]),
    ]
    gx = [0.85, 3.55, 6.25]
    for i, (title, chips) in enumerate(groups):
        add_card(s, gx[i], 5.18, 2.5, 1.45, fill=CARD_GRAY, line=RGBColor(0xC7, 0xD6, 0xE8), lw=0.8)
        add_text(s, gx[i] + 0.08, 5.27, 2.32, 0.22, title, size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        yy = 5.52
        for ch in chips:
            add_card(s, gx[i] + 0.16, yy, 2.18, 0.28, fill=WHITE, line=RGBColor(0xAF, 0xC2, 0xD9), lw=0.7)
            add_text(s, gx[i] + 0.22, yy + 0.06, 2.05, 0.14, ch, size=9.5, align=PP_ALIGN.CENTER)
            yy += 0.35

    add_bullets(
        s,
        9.66,
        1.3,
        2.95,
        2.05,
        "我们的技术主张",
        [
            "把医疗安全放在第一优先级",
            "把每次记录变成可执行干预",
            "把短期使用变成长期陪伴",
        ],
    )

    add_bullets(
        s,
        9.66,
        3.45,
        2.95,
        1.75,
        "项目阶段目标",
        [
            "立项期完成端到端Demo闭环",
            "试点期沉淀可审计规则与知识资产",
        ],
        title_color=ACCENT,
    )

    add_card(s, 9.66, 5.35, 2.95, 1.3, fill=CARD_BLUE, line=RGBColor(0xB6, 0xC8, 0xE0), lw=0.8)
    add_text(
        s,
        9.8,
        5.55,
        2.65,
        0.9,
        "价值落点：\n让患者持续记录、看懂趋势，\n并在风险出现前得到干预。",
        size=10.8,
    )

    add_footer(s, "1/4", "参考：developers.weixin.qq.com；ima.qq.com；J Med Internet Res 2020;22(9):e20283")


def slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_topbar(
        s,
        "核心技术（2/4）创新一：零负担采集，提升随访完成率",
        "将高频记录压缩到10秒级，降低失访并提升数据连续性",
    )

    add_card(s, 0.55, 1.1, 5.8, 5.0, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 6.55, 1.1, 6.2, 5.0, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 0.55, 6.2, 12.2, 0.75, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)

    add_text(s, 0.75, 1.22, 2.8, 0.24, "交互入口：2D人体点选/涂抹", size=12.5, bold=True, color=PRIMARY)

    # mock two views
    add_card(s, 1.05, 1.72, 2.15, 3.9, fill=RGBColor(0xFB, 0xFD, 0xFF), line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.8)
    add_card(s, 3.55, 1.72, 2.15, 3.9, fill=RGBColor(0xFB, 0xFD, 0xFF), line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.8)
    add_text(s, 1.65, 1.85, 0.95, 0.2, "前视图", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, 4.15, 1.85, 0.95, 0.2, "后视图", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    dots = [(1.95, 2.55), (1.95, 3.1), (1.95, 4.45), (2.4, 4.45), (1.5, 4.45), (4.45, 2.55), (4.45, 3.1), (3.95, 4.45), (4.95, 4.45)]
    for x, y in dots:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.3), Inches(0.3))
        d.fill.solid()
        d.fill.fore_color.rgb = RGBColor(0xC7, 0xE2, 0xF5)
        d.line.color.rgb = RGBColor(0x9C, 0xC2, 0xDE)

    hot1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.7), Inches(3.45), Inches(0.58), Inches(0.58))
    hot1.fill.solid(); hot1.fill.fore_color.rgb = RGBColor(0xF3, 0x9D, 0x93); hot1.line.color.rgb = ALERT
    hot2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.02), Inches(3.35), Inches(0.72), Inches(0.72))
    hot2.fill.solid(); hot2.fill.fore_color.rgb = RGBColor(0xED, 0x74, 0x74); hot2.line.color.rgb = ALERT

    add_text(s, 0.8, 5.7, 5.3, 0.35, "记录方式：点选部位 + 强度滑条 + 一键保存，不要求长文本输入。", size=10.5, color=MUTED)

    add_text(s, 6.75, 1.22, 2.8, 0.24, "AI结构化采集流程", size=12.5, bold=True, color=PRIMARY)
    flow = [
        "1) 用户一句话描述疼痛变化与用药情况",
        "2) ASR转写并抽取字段（部位/强度/诱因/疗效）",
        "3) 弹出确认卡片（确认 / 修改）",
        "4) 自动入库并触发后续评估与宣教",
    ]
    fy = 1.75
    for i, item in enumerate(flow):
        fill = CARD_BLUE if i % 2 == 0 else CARD_TEAL
        add_card(s, 6.82, fy, 5.66, 0.85, fill=fill, line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.8)
        add_text(s, 7.02, fy + 0.26, 5.3, 0.32, item, size=11.2)
        if i < len(flow) - 1:
            arrow_down(s, 9.52, fy + 0.86, 0.24, 0.17, color=ACCENT)
        fy += 1.02

    add_card(s, 6.82, 5.12, 5.66, 0.74, fill=CARD_RED, line=ALERT, lw=0.8)
    add_text(s, 7.02, 5.36, 5.3, 0.24, "触发式补采集：步数骤降 / 服药后X小时 / 连续未记录 → 主动追问1句", size=10.8)

    # bottom KPIs
    kpis = [
        ("采集耗时", "目标≤10秒", PRIMARY, CARD_BLUE),
        ("字段结构化", "自动完成", ACCENT, CARD_TEAL),
        ("失访召回", "状态机触发", ALERT, CARD_RED),
    ]
    kx = [0.72, 4.78, 8.84]
    for i, (k, v, c, f) in enumerate(kpis):
        add_card(s, kx[i], 6.34, 3.75, 0.46, fill=f, line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.8)
        add_text(s, kx[i] + 0.08, 6.48, 1.45, 0.18, k, size=10.5, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(s, kx[i] + 1.55, 6.48, 2.08, 0.18, v, size=11, bold=True, color=c, align=PP_ALIGN.CENTER)

    add_footer(s, "2/4", "参考：J Med Internet Res 2020;22(9):e20283；团队访谈与问卷调研")


def slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_topbar(
        s,
        "核心技术（3/4）创新二：分级预警与个性化推送",
        "静态阈值 + 动态状态机双路径，兼顾安全底线与干预效率",
    )

    add_card(s, 0.55, 1.1, 7.3, 5.75, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 8.05, 1.1, 4.7, 5.75, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)

    add_text(s, 0.75, 1.22, 2.8, 0.24, "双路径风险识别", size=12.5, bold=True, color=PRIMARY)

    add_card(s, 0.95, 1.72, 3.2, 2.0, fill=CARD_RED, line=ALERT, lw=0.9)
    add_bullets(
        s,
        1.08,
        1.88,
        2.92,
        1.7,
        "Path A：静态阈值/突变",
        ["NRS≥8 直接触发高危", "较7日均值突增≥3且当日≥6", "用于急性风险即时兜底"],
        title_color=ALERT,
    )

    add_card(s, 4.3, 1.72, 3.2, 2.0, fill=CARD_TEAL, line=ACCENT, lw=0.9)
    add_bullets(
        s,
        4.43,
        1.88,
        2.92,
        1.7,
        "Path B：治疗周期状态机",
        ["术后期/康复期/药物敏感期", "阈值动态调整 + 连续性判定", "识别慢性风险与依从性下降"],
        title_color=RGBColor(0x19, 0x8F, 0x84),
    )

    arrow_down(s, 2.5, 3.82, 0.25, 0.2)
    arrow_down(s, 5.85, 3.82, 0.25, 0.2)

    add_card(s, 2.3, 4.1, 3.85, 0.95, fill=CARD_BLUE, line=PRIMARY, lw=0.9)
    add_text(s, 2.45, 4.35, 3.55, 0.24, "Risk Level 1 / 2 / 3", size=13.5, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, 2.45, 4.62, 3.55, 0.2, "响应调度：通知 / 报告 / 宣教", size=11, align=PP_ALIGN.CENTER)

    add_card(s, 1.35, 5.3, 5.8, 1.32, fill=CARD_GRAY, line=RGBColor(0xC3, 0xD3, 0xE7), lw=0.8)
    add_text(
        s,
        1.6,
        5.72,
        5.3,
        0.8,
        "实时监听：触发器/云函数执行；\n审计输出：rule_id + reason + event_time，便于复盘与优化。",
        size=11,
        align=PP_ALIGN.CENTER,
    )

    add_text(s, 8.25, 1.22, 3.3, 0.24, "分级响应（系统自动执行）", size=12.5, bold=True, color=PRIMARY)

    levels = [
        ("Level 1 高危", "触发：NRS≥8或突增>3", "动作：强提醒 + 就医建议模板 + 紧急摘要", ALERT),
        ("Level 2 中危", "触发：5-7分持续≥3天", "动作：建议联系医生 + 并发症注意事项推送", RGBColor(0xD9, 0x8B, 0x16)),
        ("Level 3 依从性", "触发：超过X天未记录", "动作：温和召回 + 24h延迟队列 + 亲情号提醒", RGBColor(0x4B, 0x89, 0xD6)),
    ]
    ly = 1.72
    for name, t1, t2, c in levels:
        add_card(s, 8.25, ly, 4.3, 1.45, fill=WHITE, line=c, lw=0.9)
        add_text(s, 8.42, ly + 0.14, 2.0, 0.22, name, size=12, bold=True, color=c)
        add_text(s, 8.42, ly + 0.48, 3.95, 0.2, t1, size=10.8)
        add_text(s, 8.42, ly + 0.8, 3.95, 0.42, t2, size=10.8)
        ly += 1.7

    add_footer(s, "3/4", "参考：支瑞聪等,2020；朱南希等,2022；团队预警规则设计")


def slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_topbar(
        s,
        "核心技术（4/4）创新三：AI深度融合与可持续壁垒",
        "以“检索增强 + 脱敏治理 + 安全护栏”构建可持续迭代技术中台",
    )

    add_card(s, 0.55, 1.1, 12.2, 3.25, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)
    add_card(s, 0.55, 4.55, 12.2, 2.3, fill=WHITE, line=RGBColor(0xB9, 0xCA, 0xE1), lw=1.0)

    add_text(s, 0.75, 1.22, 3.6, 0.24, "AI闭环流程（时序）", size=12.5, bold=True, color=PRIMARY)

    flow = [
        ("用户提问", CARD_BLUE, PRIMARY),
        ("AI网关脱敏", CARD_RED, ALERT),
        ("知识检索增强\n（平台知识库 + IMA知识库）", CARD_TEAL, RGBColor(0x19, 0x90, 0x85)),
        ("安全护栏与通俗化", CARD_RED, ALERT),
        ("返回短句答案+引用", CARD_BLUE, PRIMARY),
    ]
    fx, fy = 0.95, 2.0
    fw, fh, fg = 2.2, 1.52, 0.12
    for i, (name, fill, color) in enumerate(flow):
        x = fx + i * (fw + fg)
        add_card(s, x, fy, fw, fh, fill=fill, line=color, lw=0.9)
        add_text(s, x + 0.1, fy + 0.5, fw - 0.2, 0.58, name, size=11.8, bold=True, color=color, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            arrow_right(s, x + fw + 0.02, fy + 0.63, 0.1, 0.18)

    add_card(s, 0.95, 3.65, 11.4, 0.55, fill=CARD_GRAY, line=RGBColor(0xC3, 0xD3, 0xE7), lw=0.8)
    add_text(s, 1.2, 3.83, 10.9, 0.24, "隐私原则：用户画像与病历原文不外发；外部能力仅接收脱敏后的最小必要Query。", size=10.8, color=MUTED, align=PP_ALIGN.CENTER)

    moats = [
        ("壁垒1：知识资产化", ["指南/共识/循证内容结构化沉淀", "输出带引用，可追溯可质控", "专家审核保障医学可靠性"], PRIMARY, CARD_BLUE),
        ("壁垒2：安全可控", ["规则先决策，AI做解释与教育", "高危场景使用固定就医模板", "权限、同意、审计全链路留痕"], RGBColor(0x19, 0x8F, 0x84), CARD_TEAL),
        ("壁垒3：规模化能力", ["云函数调度支撑高并发低运维", "可扩展接入更多可穿戴平台", "B2B2C路径支撑商业化推广"], ALERT, CARD_RED),
    ]

    mx = [0.85, 4.9, 8.95]
    for i, (title, lines, color, fill) in enumerate(moats):
        add_card(s, mx[i], 4.85, 3.7, 1.85, fill=fill, line=color, lw=0.9)
        add_bullets(s, mx[i] + 0.08, 4.98, 3.54, 1.65, title, lines, title_color=color)

    add_card(s, 0.85, 6.72, 11.9, 0.16, fill=CARD_BLUE, line=RGBColor(0xB8, 0xCA, 0xE1), lw=0.7)
    add_text(s, 1.2, 6.71, 11.2, 0.16, "结论：项目已形成“技术可行 + 医疗可信 + 商业可扩”的三重成立条件。", size=10.5, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    add_footer(s, "4/4", "参考：ima.qq.com；developers.weixin.qq.com；彭进亚等,2016；朱南希等,2022；李冬,2019")


def build(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)

    prs.save(output_path)


if __name__ == "__main__":
    output = "d:/project/pain_alleviation/SoothPal_核心技术_评审版_v3.pptx"
    build(output)
    print(output)
